"""MVP 제안서 양식 pptx에 마중 내용을 채워 제출본을 생성.

문체 기준: 스킬 fin-biz-tone (금융권 비즈니스 한국어, 은어 순화, em dash 금지,
완전한 구 단위 개조식, 모든 수치에 출처 또는 계산식).

실행:  cd majung/docs && python build_pptx.py
"""
from __future__ import annotations

import copy
from pathlib import Path

from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent.parent
TEMPLATE = ROOT / "[데이콘] JB금융그룹 Fin AI Challenge MVP제안서 양식.pptx"
OUTPUT = ROOT / "[제출]MVP제안서_마중.pptx"

FONT = "맑은 고딕"  # Windows 기본 설치 폰트. Pretendard 미설치 환경의 대체 렌더로 인한 줄바꿈과 겹침 방지
NAVY = RGBColor(0x1F, 0x4E, 0x79)
BLUE = RGBColor(0x2E, 0x75, 0xB6)
LBLUE = RGBColor(0xDE, 0xEB, 0xF7)
GREEN = RGBColor(0x1A, 0x7F, 0x37)
LGREEN = RGBColor(0xE6, 0xF4, 0xEA)
RED = RGBColor(0xC0, 0x39, 0x2B)
LRED = RGBColor(0xFD, 0xEC, 0xEA)
AMBER = RGBColor(0x9A, 0x67, 0x00)
GRAY = RGBColor(0x70, 0x70, 0x70)
BLACK = RGBColor(0x22, 0x22, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _ea(run):
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find(qn("a:ea"))
    if ea is None:
        ea = rpr.makeelement(qn("a:ea"), {})
        rpr.append(ea)
    ea.set("typeface", FONT)


def _run(p, text, size, color, bold, italic=False):
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size); f.bold = bold; f.italic = italic; f.color.rgb = color; f.name = FONT
    _ea(r)
    return r


def add_text(slide, l, t, w, h, blocks):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, b in enumerate(blocks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(3); p.space_before = Pt(0)
        _run(p, b[0], b[1], b[2], b[3], b[4] if len(b) > 4 else False)
    return tb


def add_box(slide, l, t, w, h, lines, fill, line, radius=True, anchor=MSO_ANCHOR.MIDDLE):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shp, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = line; sp.line.width = Pt(1.25)
    sp.shadow.inherit = False
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right"):
        setattr(tf, m, Inches(0.05))
    for m in ("margin_top", "margin_bottom"):
        setattr(tf, m, Inches(0.02))
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(0); p.space_before = Pt(0)
        _run(p, ln[0], ln[1], ln[2], ln[3])
    return sp


def add_arrow(slide, l, t, w, color=GRAY, h=0.16, down=False):
    shp = MSO_SHAPE.DOWN_ARROW if down else MSO_SHAPE.RIGHT_ARROW
    sp = slide.shapes.add_shape(shp, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp


def add_table(slide, l, t, w, h, data, col_widths, fsize=10.5):
    rows, cols = len(data), len(data[0])
    gt = slide.shapes.add_table(rows, cols, Inches(l), Inches(t), Inches(w), Inches(h)).table
    for j, cw in enumerate(col_widths):
        gt.columns[j].width = Inches(cw)
    for i, row in enumerate(data):
        for j, val in enumerate(row):
            c = gt.cell(i, j)
            c.margin_left = Inches(0.06); c.margin_right = Inches(0.06)
            c.margin_top = Inches(0.02); c.margin_bottom = Inches(0.02)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]; c.text_frame.word_wrap = True
            head = i == 0
            _run(p, val, fsize, WHITE if head else BLACK, head)
            c.fill.solid()
            c.fill.fore_color.rgb = NAVY if head else (LBLUE if i % 2 == 0 else WHITE)
    return gt


def add_quant_chart(slide, l, t, w, h):
    cd = CategoryChartData()
    cd.categories = ["보수(3%)", "기본(7%)", "낙관(12%)"]
    cd.add_series("대환 잔액(억 원)", (297, 693, 1188))
    cd.add_series("연 이자수익(억 원)", (28.5, 66.5, 113.9))
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                                Inches(l), Inches(t), Inches(w), Inches(h), cd)
    ch = gf.chart
    ch.has_title = True; ch.chart_title.text_frame.text = "전환율별 대환 잔액 및 이자수익 추정 (단위: 억 원)"
    ch.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
    ch.has_legend = True; ch.legend.position = XL_LEGEND_POSITION.BOTTOM
    ch.legend.include_in_layout = False
    ch.legend.font.size = Pt(9)
    ch.category_axis.tick_labels.font.size = Pt(11)
    ch.value_axis.tick_labels.font.size = Pt(10)
    return ch


def remove_guidance(slide):
    for sh in list(slide.shapes):
        if sh.has_text_frame and sh.text_frame.text.strip().startswith("[작성방법]"):
            sh._element.getparent().remove(sh._element)


def remove_frame_title(slide):
    """프레임 안 중제목 제거. 상단 헤더 바에 동일 제목이 있어 중복 표기가 되는 문제 해소."""
    import re
    lo, hi = Inches(1.0), Inches(1.7)
    for sh in list(slide.shapes):
        if not sh.has_text_frame or sh.top is None:
            continue
        if lo <= sh.top <= hi and re.match(r"^\d+\.\s", sh.text_frame.text.strip()):
            sh._element.getparent().remove(sh._element)


def chain(slide, items, top, h, box_w, arrow_w, start_l=None):
    n = len(items)
    total = n * box_w + (n - 1) * arrow_w
    x = start_l if start_l is not None else (13.333 - total) / 2
    ay = top + (h - 0.16) / 2
    for i, (lines, fill, line) in enumerate(items):
        add_box(slide, x, top, box_w, h, lines, fill, line)
        if i < n - 1:
            add_arrow(slide, x + box_w + 0.02, ay, arrow_w - 0.06)
        x += box_w + arrow_w


def src(slide, text):
    """슬라이드 하단 출처 각주 (콘텐츠 프레임 아래, 푸터 라인 위)."""
    add_text(slide, 0.95, 6.60, 11.5, 0.35, [(text, 8.5, GRAY, False)])


def set_cell(cell, lines):
    """양식 표 셀 텍스트 교체 (테두리와 음영 등 틀은 보존, 안내 문구 제거 효과)."""
    tf = cell.text_frame
    tf.word_wrap = True
    for para in list(tf.paragraphs[1:]):
        para._p.getparent().remove(para._p)
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)
    for i, b in enumerate(lines):
        p = p0 if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2); p.space_before = Pt(0)
        _run(p, b[0], b[1], b[2], b[3], b[4] if len(b) > 4 else False)


def fill_form_table(slide, mapping):
    """양식 표의 1열 라벨이 mapping 키와 일치하면 2열 셀을 채움."""
    for sh in slide.shapes:
        if sh.has_table:
            for row in sh.table.rows:
                label = row.cells[0].text.strip()
                if label in mapping:
                    set_cell(row.cells[1], mapping[label])


CAPTURES = Path(__file__).resolve().parent / "captures"


def add_capture_row(slide, files, top, height, caption):
    """데모 캡처 PNG를 비율 유지, 가운데 정렬로 한 줄 배치하고 아래에 공동 캡션. 전부 존재할 때만 삽입."""
    from PIL import Image
    paths = [CAPTURES / f for f in files]
    if not all(pp.exists() for pp in paths):
        return False
    widths = []
    for pp in paths:
        with Image.open(pp) as im:
            widths.append(height * im.width / im.height)
    gap = 0.4
    x = (13.333 - (sum(widths) + gap * (len(paths) - 1))) / 2
    for pp, w in zip(paths, widths):
        slide.shapes.add_picture(str(pp), Inches(x), Inches(top), height=Inches(height))
        x += w + gap
    tb = add_text(slide, 0.95, top + height + 0.05, 11.4, 0.25, [(caption, 9, GRAY, False)])
    for para in tb.text_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
    return True


def clone_slide(prs, src_idx):
    """템플릿 슬라이드를 헤더 바·푸터 로고까지 복제해 맨 뒤에 추가하고 반환.

    add_slide가 만든 기본 placeholder는 제거하고 원본 도형을 deepcopy한 뒤,
    복사된 그림(로고)의 이미지 관계(r:embed)를 새 슬라이드 관계로 재연결한다.
    """
    src = prs.slides[src_idx]
    new = prs.slides.add_slide(src.slide_layout)
    for shp in list(new.shapes):
        shp._element.getparent().remove(shp._element)
    for shp in src.shapes:
        new.shapes._spTree.append(copy.deepcopy(shp._element))
    for shp in new.shapes:
        for blip in shp._element.iter(qn("a:blip")):
            rid = blip.get(qn("r:embed"))
            if rid and rid in src.part.rels:
                img = src.part.rels[rid].target_part
                blip.set(qn("r:embed"), new.part.relate_to(img, RT.IMAGE))
    return new


def fill_diagram_slide(slide, header, title, img_name):
    """복제한 부록 슬라이드에 헤더 교체 + 슬라이드 제목 + 다이어그램을 푸터와
    충돌 없이 최대 크기로 가운데 배치한다. 비율 유지."""
    from PIL import Image
    remove_guidance(slide)
    remove_frame_title(slide)
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip().startswith("예선 산출물"):
            p0 = sh.text_frame.paragraphs[0]
            if p0.runs:
                p0.runs[0].text = header
                for r in list(p0.runs[1:]):
                    r.text = ""
            break
    add_text(slide, 0.95, 1.18, 11.4, 0.4, [(title, 14, NAVY, True)])
    path = CAPTURES / img_name
    if not path.exists():
        return
    with Image.open(path) as im:
        iw, ih = im.size
    max_w, max_h = 12.4, 5.05            # 인치 (top 1.80~6.85, 푸터 7.02 위)
    w = max_w; h = w * ih / iw
    if h > max_h:
        h = max_h; w = h * iw / ih
    left = (13.333 - w) / 2
    slide.shapes.add_picture(str(path), Inches(left), Inches(1.80),
                             width=Inches(w), height=Inches(h))


def main():
    p = Presentation(str(TEMPLATE))
    s = p.slides
    # 부록: 다이어그램 전용 슬라이드 3장(헤더·푸터 포함 복제). 본문 채우기 전에 깨끗한 복제 확보
    diag = [clone_slide(p, 3) for _ in range(3)]

    # ════════ Slide 0 · 표지 ════════
    tb = add_text(s[0], 2.9, 4.0, 7.5, 0.8, [
        ("마중", 18, NAVY, True),
        ("외국인 근로자 맞춤형 위임 뱅킹 에이전트 · 전북은행·브라보코리아·한패스 연계", 11, BLACK, False),
    ])
    for para in tb.text_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
    fill_form_table(s[0], {
        "팀명": [("마중", 11, BLACK, True)],
        "서비스명": [("마중", 11, NAVY, True)],
    })

    # ════════ Slide 1 · Summary ════════
    remove_guidance(s[1])
    fill_form_table(s[1], {
        "팀명": [("마중", 11, BLACK, False)],
        "팀원": [("문기훈 (1인 참가)", 11, BLACK, False)],
        "서비스명": [
            ("마중", 13, NAVY, True),
            ("외국인 근로자 특화 위임형 뱅킹 에이전트", 10.5, BLACK, False),
            ("'마중': 낯선 땅에 도착한 손님을 먼저 나가 맞이한다는 우리말", 9, GRAY, False),
        ],
        "주제": [
            ("자유주제: 외국인 근로자가 신용 사다리를 오르도록 돕는 뱅킹 에이전트 (전북은행·브라보코리아·한패스 연계)",
             11, BLACK, False),
        ],
        "한줄소개": [
            ('"Lương về thì gửi cho mẹ, khi tỷ giá tốt, tối đa 2 triệu."', 12.5, NAVY, True),
            ("(월급이 입금되면 환율이 유리할 때 어머니께 최대 200만 원까지 송금해 줘)", 9.5, GRAY, False),
            ("모국어 한마디로 송금을 위임하고, 고금리 사채를 전북은행 대출로 갈아타도록 안내하여 우량 고객으로 키우는 AI 에이전트",
             10.5, BLACK, False),
        ],
    })

    # ════════ Slide 2 · 문제 정의 ════════
    remove_guidance(s[2]); remove_frame_title(s[2])
    add_text(s[2], 0.95, 1.30, 11.4, 0.4,
             [("입국 과정에서 떠안은 고금리 사채가 외국인의 제도권 금융 진입을 원천 차단", 13.5, NAVY, True)])
    journey = [
        ([("브로커 사채 이용", 10.5, RED, True), ("원금 1,500만·연 20~40%", 8, BLACK, False)], LRED, RED),
        ([("과도한 이자 부담", 10.5, RED, True), ("급여 상당액이 이자로 증발", 8, BLACK, False)], LRED, RED),
        ([("불법 추심 및 협박", 10.5, RED, True), ("연체 시 강제 출국 협박", 8, BLACK, False)], LRED, RED),
        ([("제도권 금융 거절", 10.5, RED, True), ("금융 이력 부족으로 반려", 8, BLACK, False)], LRED, RED),
        ([("사채의 늪 재유입", 10.5, RED, True)], LRED, RED),
    ]
    chain(s[2], journey, top=2.16, h=0.8, box_w=2.0, arrow_w=0.42)
    add_text(s[2], 0.95, 3.02, 11.4, 1.05, [
        ("[현황] 금융 사각지대에 방치된 외국인 근로자", 11, NAVY, True),
        ("· 5대 시중은행 內 외국인 전용 신용대출 상품 부재 (출처: 머니투데이)", 10.5, BLACK, False),
        ("· 부산 적발 사례: 외국인 9,120명 대상 162억 원 대출, 연 최고 154% 폭리 (출처: 경향신문, 2025.11)", 10.5, BLACK, False),
        ("· 제도권의 대출 거절이 사채 의존도를 높이는 악순환 구조 고착화", 10.5, BLACK, False),
    ])
    add_text(s[2], 0.95, 4.28, 11.4, 1.1, [
        ("[위협] 경쟁행의 공격적인 급여계좌 선점 가속화", 11, NAVY, True),
        ("· [신한] '급여 3개월 수령'을 대출 필수 조건화하여 초기 계좌 이탈 방지 (2025.9)", 10.5, BLACK, False),
        ("· [전북] 외국인 신용대출 점유율 72%, 누적 고객 24만 명 달성 (출처: 머니투데이, 2025.4)", 10.5, BLACK, False),
        ("· 시사점: 선제적 이탈 방지 및 방어 전략 부재 시 핵심 우량 고객 대규모 이탈 불가피", 10.5, BLACK, False),
    ])
    add_text(s[2], 0.95, 5.50, 11.4, 0.95, [
        ("[전략] 신규 고객 유치 출혈경쟁을 탈피한 '대환 중심의 질적 성장'", 11, NAVY, True),
        ("· 핵심 KPI: 고객 생애 가치 제고 및 타행/사채 대환 전환율", 10.5, BLACK, False),
        ("· 배경: 외국인의 최초 거래은행은 사업주가 결정하므로, 입사 이후 대환을 통한 주거래 유도가 가장 현실적인 전략", 10.5, BLACK, False),
    ])

    # ════════ Slide 3 · 솔루션 개요 ════════
    remove_guidance(s[3]); remove_frame_title(s[3])
    add_text(s[3], 0.95, 1.30, 11.4, 0.4,
             [("유입은 자동 송금으로, 핵심 수익은 사채 대환에서 창출", 14, NAVY, True)])
    add_box(s[3], 0.95, 2.2, 5.55, 0.95, [
        ("Step 1. 위임 송금을 통한 고객 확보 (유입 채널)", 11, NAVY, True),
        ("급여 입금 시 지정된 환율 조건에 맞춰 자동 송금 실행 및 이상 거래 보류", 9.5, BLACK, False)], LBLUE, BLUE)
    add_box(s[3], 6.80, 2.2, 5.55, 0.95, [
        ("Step 2. 사채 대환을 통한 신용 형성 (수익 창출)", 11, GREEN, True),
        ("고금리 사채(연 30%)를 전북은행 대출(연 13.59%)로 대환 유도 및 적금 연계", 9.5, BLACK, False)], LGREEN, GREEN)
    add_box(s[3], 0.85, 3.5, 1.75, 1.05,
            [("요청 및 급여 입금", 10, BLACK, True), ("작동 시작 시점", 9, BLACK, False)], RGBColor(0xF2, 0xF2, 0xF2), GRAY)
    add_arrow(s[3], 2.62, 3.95, 0.33)
    add_box(s[3], 2.97, 3.5, 1.95, 1.05,
            [("LLM 의도 해석", 10, BLACK, True), ("자연어 해석 전담", 8.5, BLACK, False),
             ("실행 권한 없음", 8.5, BLACK, False)],
            RGBColor(0xF2, 0xF2, 0xF2), GRAY)
    add_arrow(s[3], 4.94, 3.95, 0.33)
    add_box(s[3], 5.29, 3.5, 4.55, 1.08, [("3중 보안 게이트 (규칙 기반 코드가 판정)", 9.5, NAVY, True)],
            LBLUE, NAVY, anchor=MSO_ANCHOR.TOP)
    for gx, ttl, det in [(5.42, "게이트 A", "위임 검증"), (6.85, "게이트 B", "한도·환율"), (8.28, "게이트 C", "수취인·AML")]:
        add_box(s[3], gx, 3.97, 1.42, 0.52, [(ttl, 9, WHITE, True), (det, 8, WHITE, False)], BLUE, NAVY)
    add_arrow(s[3], 9.86, 3.95, 0.33)
    add_box(s[3], 10.21, 3.5, 2.28, 1.05,
            [("최종 실행", 10, GREEN, True), ("모국어 통지 및 즉시 철회 지원", 8.5, BLACK, False)], LGREEN, GREEN)
    add_text(s[3], 0.95, 4.8, 11.4, 1.6, [
        ("송금 실행은 LLM이 아니라 정해진 규칙대로만 작동하는 코드가 담당하여 금융 사고를 차단", 12, NAVY, True),
        ("· LLM은 고객의 자연어 의도 해석 및 번역만 수행하며, 실제 송금액 산정과 실행 여부는 내부 '정해진 규칙대로만 실행되어 결과가 늘 같은 규칙 기반 코드'가 통제", 11, BLACK, False),
        ("브라보코리아 앱에 내장하는 대환 전환 엔진", 12, NAVY, True),
        ("· 기존 휴면 앱 다운로드 20만 건과 외국인 고객 24만 명을 대환대출·적금으로 재활성화", 11, BLACK, False),
    ])

    # ════════ Slide 4 · 주요 기능 및 아키텍처 ════════
    remove_guidance(s[4]); remove_frame_title(s[4])
    add_text(s[4], 0.95, 1.30, 11.4, 0.4,
             [("사용자 페인포인트를 해결하는 3중 게이트 기반의 안전한 자금 통제", 13.5, NAVY, True)])
    add_table(s[4], 0.95, 2.12, 11.4, 1.5, [
        ["외국인 근로자 핵심 과업", "마중 솔루션"],
        ["\"환율이 유리할 때 송금하고 싶지만 일하느라 매번 확인하기 어렵다\"", "급여 입금 감지 후 목표 환율 도달 시 자동 송금 실행"],
        ["\"잘못 송금되거나 사기를 당할까 봐 두렵다\"", "3중 게이트 기반 자동 보류 및 자금세탁방지 선차단"],
        ["\"사채 이자가 벅차지만, 은행 대출 방법을 모른다\"", "대환 가심사 후 전북은행 심사엔진 회부 (연 약 250만 원 이자 절감)"],
        ["\"한국어 금융 서류를 읽고 이해하기 불가능하다\"", "위임장 모국어 렌더링, 건별 모국어 푸시 통지 및 철회 기능 제공"],
    ], col_widths=[5.6, 5.8], fsize=10.5)
    add_text(s[4], 0.95, 3.95, 11.4, 0.3, [("Agent 루프 아키텍처: 지속적인 검증 및 한도 최적화", 11, NAVY, True)])
    loop = [
        ([("판단", 11, WHITE, True), ("게이트 A·B·C 검증", 8.5, WHITE, False)], NAVY, NAVY),
        ([("행동", 11, WHITE, True), ("실행 / 보류 / 모국어 알림", 8.5, WHITE, False)], BLUE, NAVY),
        ([("검증/개선", 11, WHITE, True), ("금융 이력 누적 → 위임 한도 상향", 8.5, WHITE, False)], GREEN, GREEN),
    ]
    chain(s[4], loop, top=4.28, h=0.7, box_w=3.0, arrow_w=0.45, start_l=0.95)
    add_text(s[4], 0.95, 5.16, 11.4, 1.3, [
        ("▶ 세부 처리 단계 (평가 6단계): 수집 → 검색 → 판단 → 생성 → 검증 → 후속액션", 11, NAVY, True),
        ("· 기능 구현 완료: 위임장 발급, 조건부 자동 송금, 의심 거래 보류, 사기 거래 차단, 대환대출 회부 (단위 테스트 20건 전원 통과)", 11, BLACK, False),
        ("· 시스템 설계 경계: 대출 최종 승인 권한은 전북은행 고유 엔진에 위임, 단일 사고 발생 시 최대 손실액을 위임 한도 내로 제한",
         11, BLACK, False),
    ])

    # ════════ Slide 5 · 데이터 및 기술 ════════
    remove_guidance(s[5]); remove_frame_title(s[5])
    add_text(s[5], 0.95, 1.30, 11.4, 0.4,
             [("개인정보 비식별 처리와 허용형 오픈소스 기반의 규제 준수 아키텍처", 13.5, NAVY, True)])
    add_text(s[5], 0.95, 2.18, 11.4, 1.5, [
        ("· LLM: 교체 가능한 연동 구조로 상용 LLM(Claude 등)을 붙이도록 설계, 의도 해석·번역 전담. 예선 PoC는 규칙 기반 대체 방식(본선에서 코드 변경 없이 교체), RAG 미적용", 11, BLACK, False),
        ("  계좌·금액·실명은 미전송하고 비식별 의도 텍스트만 전달, 판정은 SQLite 규칙 기반 코드 전담. 상용화 시 국내 리전 또는 행내 sLLM 전환 가능", 10, GRAY, False),
        ("· 프라이버시 보호: 외부 LLM 호출 시 민감 식별정보를 비식별 치환 처리, 판정 및 감사 로그는 내부 망에서 단독 처리", 11, BLACK, False),
        ("· 데이터: 급여·송금·사채 데이터 전부 합성(실 개인정보 사용 배제) / 스택: Python·FastAPI·SQLite·React·MCP SDK (전부 MIT·BSD 등 허용형 라이선스)", 11, BLACK, False),
    ])
    add_text(s[5], 0.95, 3.78, 11.4, 0.32, [("[핵심 기술 제약 사항 및 선제적 대응 방안]", 12, NAVY, True)])
    add_table(s[5], 0.95, 4.12, 11.4, 1.95, [
        ["예상되는 기술적 제약", "마중의 대응 전략"],
        ["코어뱅킹 실연동 제약 (망분리 및 보안성 심의)", "모의 API 기반 PoC 선행 구현, 실제 코어 연동은 서비스 상용화 단계로 분리 설계"],
        ["LLM 환각에 따른 금융 사고 책임", "LLM의 송금 실행 권한 박탈, 허용 수취인 명단 관리 및 1회 송금 한도 상한 설정"],
        ["환율 예측의 불확실성", "단순 예측 배제, '7일 이동평균 대비 1% 이상 시' 등 명시적 규칙 적용 및 환율 데이터 60초 지연 시 거래 정지"],
        ["외국인 대상 자동화 송금에 따른 자금세탁 리스크", "게이트 C 내부 점수화 적용, 이상 징후 시 의심거래보고 큐 등록 및 송금 선차단"],
    ], col_widths=[4.4, 7.0], fsize=10.5)

    # ════════ Slide 6 · 사용자 데모 시나리오 ════════
    remove_guidance(s[6]); remove_frame_title(s[6])
    add_text(s[6], 0.95, 1.30, 11.4, 0.4,
             [("입국부터 대환까지 이어지는 전 과정 시나리오 검증 완료", 13.5, NAVY, True)])
    add_text(s[6], 0.95, 1.78, 11.4, 0.35,
             [("테스트 페르소나: 한국어에 미숙하며, 입국 자금 마련을 위해 브로커 사채 1,500만 원을 이용 중인 베트남 국적 E-9 근로자", 11, BLACK, False)])
    demo = [
        ([("① 위임장 설정", 10, WHITE, True), ("모국어 내용 확인 및 전자서명", 8.5, WHITE, False)], NAVY, NAVY),
        ([("② 조건부 송금", 10, WHITE, True), ("급여 입금 직후 목표 환율 도달 시", 8.5, WHITE, False)], GREEN, GREEN),
        ([("③ 이상 거래 보류", 10, WHITE, True), ("패턴 이탈 감지 시 일시 정지", 8.5, WHITE, False)], AMBER, AMBER),
        ([("④ 보이스피싱 차단", 10, WHITE, True), ("비정상 OTP 요구 시 즉각 차단", 8.5, WHITE, False)], RED, RED),
        ([("⑤ 사채 대환 제안", 10, WHITE, True), ("신용 데이터 기반 대출 연계", 8.5, WHITE, False)], NAVY, NAVY),
    ]
    chain(s[6], demo, top=2.28, h=1.2, box_w=2.1, arrow_w=0.28, start_l=0.85)
    add_text(s[6], 0.95, 3.66, 11.4, 0.8, [
        ("· 시나리오 ②: 고객 개입 없이 백그라운드에서 조건 충족 시 즉각 실행", 11, BLACK, False),
        ("· 시나리오 ③: 보류 내역은 모국어로 고객에게 즉시 고지되며, 내부 관리자 STR 심사 대기열로 자동 이관", 11, BLACK, False),
        ("· 시나리오 ⑤: 예상 이자 절감액과 월 상환액 등 금융소비자보호법을 준수한 모의 고지 진행 (최종 승인은 전북은행 심사로 이관)", 11, BLACK, False),
    ])
    inserted = add_capture_row(
        s[6], ["02_auto_remit.png", "03_str_hold.png", "05_refi_offer.png"],
        top=4.50, height=1.70,
        caption="실동작 화면: ② 조건부 자동 송금(철회 보장) / ③ 이상 거래 보류(점수 75) / ⑤ 대환대출 제안(투명한 정보 고지)")
    if not inserted:
        add_box(s[6], 2.4, 4.70, 8.5, 1.4, [
            ("실동작 화면 캡처 3종 삽입 예정 (② 조건부 자동 송금, ③ 이상 거래 보류, ⑤ 대환대출 제안)", 10.5, GRAY, False)],
            WHITE, GRAY)

    # ════════ Slide 7 · 기대 효과 및 비즈니스 타당성 ════════
    remove_guidance(s[7]); remove_frame_title(s[7])
    add_text(s[7], 0.95, 1.30, 11.4, 0.4,
             [("대환 전환율 7% 달성 시, 연간 이자수익 66.5억 원 창출 전망 (자체 추정치)", 13.5, NAVY, True)])
    add_quant_chart(s[7], 0.8, 1.88, 5.5, 2.78)
    add_text(s[7], 6.55, 1.88, 6.05, 2.85, [
        ("[수익기회] 대환 잔액 693억 원 = 타겟 6.6만 명 × 전환율 7% × 인당 1,500만 (자체 추정)", 11, BLACK, False),
        ("[수익기회] 연 이자수익 66.5억 원 = 693억 × 금리 스프레드 9.59%p (대손충당금·운영비 차감 전)", 11, BLACK, False),
        ("[고객가치] 1인당 연간 약 250만 원 이자 비용 절감", 11, BLACK, False),
        ("[업무효율] 대면 모집 비용 없이 고객 획득 비용 최소화, 플랫폼 연계 수수료만 발생", 11, BLACK, False),
        ("[리스크감소] AML 자동 탐지(점수 40·70 이중 임계치)와 모국어 위임장 재확인으로 불완전판매·STR 누락 위험 제거", 11, BLACK, False),
        ("[전략 연계] JB금융그룹 AI 전환 및 외국인·전략 여신 확대 방침과 직접 연계, 급여 이체를 전북은행에 모아 주거래 고객으로 정착", 11, NAVY, True),
    ])
    add_table(s[7], 0.95, 4.78, 11.4, 1.2, [
        ["유관 부서", "예상되는 내부 허들 및 우려 사항", "마중의 전략적 방어 논리"],
        ["외국인금융부", "브라보 앱 트래픽을 잠식하는 별도 서비스 아닌가?", "브라보 앱에 내장한 전환 엔진으로 기획, 앱 활성도와 LTV 동반 상승"],
        ["디지털전략부", "외부 빅테크 채널에 종속될 위험은 없는가?", "자사 앱 자체 대화형 UI가 기본, 외부 채널은 선택적 확장 수단"],
        ["여신심사부", "자동 대환에 따른 불완전판매 리스크는?", "가심사와 상품 안내까지만 수행, 최종 심사와 승인은 전북은행 고유 엔진이 전담"],
        ["준법감시부", "자연어 위임이 규제와 법적 효력을 충족하는가?", "건별 모국어 통지와 100% 철회권, AML 선차단 (PoC는 SHA-256 무결성, 상용화 시 전자서명법 적격 수단 적용·유권해석)"],
    ], col_widths=[1.7, 3.4, 6.3], fsize=9)
    add_text(s[7], 0.95, 6.06, 11.4, 0.3, [
        ("[향후 확장성] 17개국 다국어 지원, 음성 인식 고도화, 시니어/소상공인 대상 서비스 확장, 코어뱅킹 SI 실연동", 9.5, GRAY, False),
    ])
    src(s[7], "출처: 전북은행 점유율 72% 및 1.3조 목표치(JB금융지주 IR 컨퍼런스콜, 2025) / 대환잔액 및 수익 지표는 보수적 가정을 적용한 자체 추정치")

    # ════════ 부록 슬라이드: 다이어그램 3종(크게) ════════
    fill_diagram_slide(
        diag[0], "예선 산출물 – 부록 1. 시스템 아키텍처",
        "시스템 아키텍처: LLM은 의도 해석만 담당하고, 실행과 판단은 규칙 기반 코드의 3중 게이트가 통제",
        "architecture.png")
    fill_diagram_slide(
        diag[1], "예선 산출물 – 부록 2. 주요 기능 흐름도",
        "주요 기능 흐름도: 위임 설정 → 급여 입금 감지 → 3중 게이트 판정 → 통과·보류·차단 분기 → 사후 처리, 사채 대환 회부",
        "flow.png")
    fill_diagram_slide(
        diag[2], "예선 산출물 – 부록 3. 유스케이스",
        "유스케이스: 외국인 근로자·관리자·전북은행 심사엔진·제휴처 간 상호작용과 시스템 경계",
        "usecase.png")

    p.save(str(OUTPUT))
    print(f"[build] 생성 완료: {OUTPUT.name}  (슬라이드 {len(s)}개)")


if __name__ == "__main__":
    main()
